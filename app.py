from flask import Flask, request, render_template, send_file, redirect, url_for, session
from pptx import Presentation
from pptx.util import Inches
import groq
import os
import io
from dotenv import load_dotenv

app = Flask(__name__)
app.secret_key = 'supersecretkey'  # Required for session management

# Load environment variables from .env file
load_dotenv()

# Set your Groq API key
groq_api_key = os.getenv("GROQ_API_KEY")

# Initialize Groq client
client = groq.Groq(api_key=groq_api_key)

def clean_text(text):
    """Remove unwanted characters from text."""
    return text.replace('\r', '').replace('_x000D_', '').strip()

def create_slide(prs, title, content, bullets=None, image_path=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = clean_text(title)

    content_placeholder = slide.placeholders[1]
    content_placeholder.text = clean_text(content)

    if bullets:
        for bullet in bullets:
            p = content_placeholder.text_frame.add_paragraph()
            p.text = clean_text(bullet)
            p.level = 1
    
    if image_path:
        slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(5))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_slides', methods=['POST'])
def generate_slides():
    presentation_description = request.form.get('presentation_description', '')
    informational_content = request.form.get('informational_content', '')
    
    # Create a hardcoded prompt with a clear example
    prompt = (
        f"Create slides based on the following description: {presentation_description}.\n\n"
        f"Here is the informational content to use in the slides: {informational_content}.\n\n"
        "Please format the slides in markdown format with clear titles, subtitles, and bullet points for each slide. "
        "Here is an example format:\n\n"
        "**Slide 1: Introduction**\n"
        "==============================\n"
        "* Title: \"Introduction to Topic\"\n"
        "* Subtitle: \"Overview of the subject\"\n"
        "* Bullet points:\n"
        "  + Key point 1\n"
        "  + Key point 2\n"
        "  + Key point 3\n"
    )
    
    # Call Groq API to generate slide content
    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama3-8b-8192",
    )
    
    # Print the raw response for troubleshooting
    print("Raw response from Groq:", flush=True)
    print(chat_completion, flush=True)
    
    # Extract the content from the response
    slide_content = chat_completion.choices[0].message.content.strip()  # Adjust according to actual response structure
    
    # Log the entire slide content for debugging purposes
    print("Slide content from Groq:", flush=True)
    print(slide_content, flush=True)

    # Parse slide content
    slides = []
    slide_sections = slide_content.split("\n\n")

    for section in slide_sections:
        section = section.strip()
        if section.startswith("**Slide"):
            lines = section.split("\n")
            title = None
            subtitle = None
            bullets = []
            for line in lines:
                line = clean_text(line.strip())
                if line.startswith("* Title:"):
                    title = line.split(":")[1].strip().strip('"')
                elif line.startswith("* Subtitle:"):
                    subtitle = line.split(":")[1].strip().strip('"')
                elif line.startswith("+"):
                    bullets.append(line[1:].strip())
            content = (subtitle if subtitle else "") + "\n" + "\n".join(bullets)
            slides.append({
                'title': title,
                'content': content,
                'bullets': bullets
            })
    
    # Store slide data in session
    session['slides'] = slides

    return redirect(url_for('preview_slides'))

@app.route('/preview_slides')
def preview_slides():
    slides = session.get('slides', [])
    return render_template('preview.html', slides=slides)

@app.route('/update_slides', methods=['POST'])
def update_slides():
    slides = session.get('slides', [])
    updated_slides = []
    
    for i in range(len(slides)):
        title = request.form.get(f'title{i+1}', slides[i]['title'])
        content = request.form.get(f'content{i+1}', slides[i]['content'])
        updated_slides.append({'title': title.strip(), 'content': content.strip()})
    
    # Update the session with the new slides data
    session['slides'] = updated_slides
    
    return redirect(url_for('preview_slides'))

@app.route('/export_slides')
def export_slides():
    slides = session.get('slides', [])
    prs = Presentation()

    for slide in slides:
        bullets = slide['content'].split("\n")[1:]  # Separate bullets from content
        create_slide(prs, slide['title'], slide['content'].split("\n")[0], bullets)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(output, download_name='presentation.pptx', as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)